from __future__ import annotations

import io
import re
from typing import Any


def safe_filename(value: str, fallback: str = "artifact") -> str:
    cleaned = re.sub(r"[^A-Za-z0-9_-]+", "_", value.strip()).strip("_")
    return (cleaned or fallback)[:80]


def render_document_docx(content: dict[str, Any]) -> bytes:
    from docx import Document
    from docx.enum.text import WD_ALIGN_PARAGRAPH
    from docx.shared import Inches, Pt, RGBColor

    doc = Document()
    section = doc.sections[0]
    section.top_margin = Inches(0.72)
    section.bottom_margin = Inches(0.72)
    section.left_margin = Inches(0.78)
    section.right_margin = Inches(0.78)
    normal = doc.styles["Normal"]
    normal.font.name = "Aptos"
    normal.font.size = Pt(10.5)
    normal.font.color.rgb = RGBColor(38, 45, 58)
    normal.paragraph_format.space_after = Pt(6)

    title = doc.add_paragraph()
    title.alignment = WD_ALIGN_PARAGRAPH.LEFT
    run = title.add_run(str(content.get("title") or "Professional document"))
    run.bold = True
    run.font.name = "Aptos Display"
    run.font.size = Pt(24)
    run.font.color.rgb = RGBColor(20, 43, 73)
    subtitle_text = str(content.get("subtitle") or "").strip()
    if subtitle_text:
        subtitle = doc.add_paragraph()
        subtitle.paragraph_format.space_after = Pt(14)
        run = subtitle.add_run(subtitle_text)
        run.italic = True
        run.font.size = Pt(10.5)
        run.font.color.rgb = RGBColor(74, 92, 112)

    for raw in content.get("sections", []):
        if not isinstance(raw, dict):
            continue
        heading = str(raw.get("heading") or "").strip()
        if heading:
            paragraph = doc.add_paragraph()
            paragraph.paragraph_format.space_before = Pt(10)
            paragraph.paragraph_format.space_after = Pt(4)
            run = paragraph.add_run(heading)
            run.bold = True
            run.font.size = Pt(13)
            run.font.color.rgb = RGBColor(21, 94, 117)
        for text in raw.get("paragraphs", []):
            doc.add_paragraph(str(text))
        for text in raw.get("bullets", []):
            doc.add_paragraph(str(text), style="List Bullet")
    closing = str(content.get("closing") or "").strip()
    if closing:
        paragraph = doc.add_paragraph(closing)
        paragraph.paragraph_format.space_before = Pt(10)
    buffer = io.BytesIO()
    doc.save(buffer)
    return buffer.getvalue()


def render_document_pdf(content: dict[str, Any]) -> bytes:
    from xml.sax.saxutils import escape

    from reportlab.lib import colors
    from reportlab.lib.pagesizes import A4
    from reportlab.lib.styles import ParagraphStyle, getSampleStyleSheet
    from reportlab.lib.units import mm
    from reportlab.platypus import Paragraph, SimpleDocTemplate, Spacer

    buffer = io.BytesIO()
    doc = SimpleDocTemplate(buffer, pagesize=A4, leftMargin=18 * mm, rightMargin=18 * mm, topMargin=18 * mm, bottomMargin=18 * mm)
    styles = getSampleStyleSheet()
    title_style = ParagraphStyle("ArtifactTitle", parent=styles["Title"], fontName="Helvetica-Bold", fontSize=22, leading=26, alignment=0, textColor=colors.HexColor("#142B49"), spaceAfter=8)
    subtitle_style = ParagraphStyle("ArtifactSubtitle", parent=styles["BodyText"], textColor=colors.HexColor("#53677D"), fontSize=10, leading=14, spaceAfter=14)
    heading_style = ParagraphStyle("ArtifactHeading", parent=styles["Heading2"], textColor=colors.HexColor("#155E75"), fontSize=13, leading=16, spaceBefore=9, spaceAfter=4)
    body_style = ParagraphStyle("ArtifactBody", parent=styles["BodyText"], fontSize=10, leading=14, spaceAfter=6)
    bullet_style = ParagraphStyle("ArtifactBullet", parent=body_style, leftIndent=12, firstLineIndent=-8, bulletIndent=2)
    story = [Paragraph(escape(str(content.get("title") or "Professional document")), title_style)]
    if content.get("subtitle"):
        story.append(Paragraph(escape(str(content["subtitle"])), subtitle_style))
    for raw in content.get("sections", []):
        if not isinstance(raw, dict):
            continue
        if raw.get("heading"):
            story.append(Paragraph(escape(str(raw["heading"])), heading_style))
        for text in raw.get("paragraphs", []):
            story.append(Paragraph(escape(str(text)), body_style))
        for text in raw.get("bullets", []):
            story.append(Paragraph("• " + escape(str(text)), bullet_style))
    if content.get("closing"):
        story.extend([Spacer(1, 8), Paragraph(escape(str(content["closing"])), body_style)])
    doc.build(story)
    return buffer.getvalue()


def render_presentation_pptx(content: dict[str, Any]) -> bytes:
    from pptx import Presentation
    from pptx.dml.color import RGBColor
    from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
    from pptx.enum.shapes import MSO_SHAPE
    from pptx.util import Inches, Pt

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]
    navy = RGBColor(16, 32, 52)
    teal = RGBColor(17, 129, 141)
    ink = RGBColor(31, 42, 55)
    muted = RGBColor(88, 105, 123)
    paper = RGBColor(247, 249, 250)

    def background(slide, color):
        fill = slide.background.fill
        fill.solid()
        fill.fore_color.rgb = color

    def textbox(slide, x, y, w, h, text, size, color, bold=False, font="Aptos", align=PP_ALIGN.LEFT):
        shape = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
        frame = shape.text_frame
        frame.clear()
        frame.word_wrap = True
        frame.vertical_anchor = MSO_ANCHOR.TOP
        paragraph = frame.paragraphs[0]
        paragraph.alignment = align
        run = paragraph.add_run()
        run.text = text
        run.font.name = font
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.color.rgb = color
        return shape

    title_slide = prs.slides.add_slide(blank)
    background(title_slide, navy)
    accent = title_slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.7), Inches(0.75), Inches(0.12), Inches(5.9))
    accent.fill.solid()
    accent.fill.fore_color.rgb = teal
    accent.line.fill.background()
    textbox(title_slide, 1.15, 1.45, 10.8, 1.7, str(content.get("title") or "Presentation"), 34, RGBColor(255, 255, 255), True, "Aptos Display")
    textbox(title_slide, 1.18, 3.25, 9.8, 1.1, str(content.get("subtitle") or "Prepared with ApplyEngine"), 17, RGBColor(184, 205, 216))
    textbox(title_slide, 1.18, 6.55, 4, 0.35, "APPLYENGINE · SKILLS", 9, RGBColor(120, 171, 180), True)

    for index, raw in enumerate(content.get("slides", []), start=1):
        if not isinstance(raw, dict):
            continue
        slide = prs.slides.add_slide(blank)
        background(slide, paper)
        bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(0.14), prs.slide_height)
        bar.fill.solid()
        bar.fill.fore_color.rgb = teal
        bar.line.fill.background()
        textbox(slide, 0.75, 0.55, 10.5, 0.35, str(raw.get("kicker") or f"{index:02d}").upper(), 10, teal, True)
        textbox(slide, 0.75, 1.0, 11.5, 0.75, str(raw.get("title") or "Key point"), 27, navy, True, "Aptos Display")
        body = str(raw.get("body") or "").strip()
        if body:
            textbox(slide, 0.78, 1.95, 11.2, 0.95, body, 15, muted)
        bullets = [str(item) for item in raw.get("bullets", []) if str(item).strip()]
        if bullets:
            shape = slide.shapes.add_textbox(Inches(0.88), Inches(3.0 if body else 2.05), Inches(11.15), Inches(3.5))
            frame = shape.text_frame
            frame.clear()
            frame.word_wrap = True
            for bullet_index, text in enumerate(bullets[:7]):
                paragraph = frame.paragraphs[0] if bullet_index == 0 else frame.add_paragraph()
                paragraph.text = "•  " + text
                paragraph.font.name = "Aptos"
                paragraph.font.size = Pt(17)
                paragraph.font.color.rgb = ink
                paragraph.space_after = Pt(13)
        textbox(slide, 11.8, 6.85, 0.65, 0.25, f"{index:02d}", 9, muted, align=PP_ALIGN.RIGHT)
        notes = str(raw.get("speaker_notes") or "").strip()
        if notes:
            try:
                slide.notes_slide.notes_text_frame.text = notes
            except Exception:
                pass
    buffer = io.BytesIO()
    prs.save(buffer)
    return buffer.getvalue()
